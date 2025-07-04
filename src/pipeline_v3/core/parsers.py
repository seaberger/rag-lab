# In the refactored datasheet_ingest_pipeline.py
import hashlib
import json
from enum import Enum
from pathlib import Path
from typing import Any

# Use absolute imports to avoid relative import issues
try:
    from storage.cache import CacheManager

    from utils.common_utils import logger
    from utils.config import PipelineConfig
    from utils.enhanced_retry import enhanced_retry_api_call
    from utils.openai_client import create_vision_client
    from utils.page_range import (
        PageProgressMonitor,
        PageRangeParser,
        get_page_count_from_pdf,
    )
except ImportError:
    # Fallback for when running from different directory
    import sys

    sys.path.append(str(Path(__file__).parent.parent))
    from storage.cache import CacheManager

    from utils.common_utils import logger
    from utils.config import PipelineConfig
    from utils.enhanced_retry import enhanced_retry_api_call
    from utils.openai_client import create_vision_client
    from utils.page_range import (
        PageProgressMonitor,
        PageRangeParser,
        get_page_count_from_pdf,
    )


def _find_poppler() -> str | None:
    """Return directory that contains pdfinfo/pdftoppm (Poppler) or None."""
    import shutil

    exe = shutil.which("pdfinfo")
    return None if exe is None else str(Path(exe).parent)


def _pdf_to_data_uris(
    pdf_path: Path,
    dpi: int = 150,
    poppler_path: str | None = None,
    page_range: str | None = None,
) -> tuple[list[str], int]:
    """Convert PDF pages to base64 data URIs for OpenAI Vision API.

    Args:
        pdf_path: Path to PDF file
        dpi: DPI for image conversion
        poppler_path: Path to Poppler binaries
        page_range: Page range specification (e.g., "1-5", "1,3,7")

    Returns:
        Tuple of (data_uris, total_page_count)
        Note: total_page_count is the original document page count,
              len(data_uris) is the number of processed pages
    """
    import base64
    import io

    from pdf2image import convert_from_path

    # Auto-discover Poppler if not provided
    if poppler_path is None:
        poppler_path = _find_poppler()
        if poppler_path is None:
            logger.warning("Poppler not found in PATH. PDF conversion may fail.")

    try:
        # First, get total page count for validation and progress tracking
        total_page_count = get_page_count_from_pdf(pdf_path)

        # Parse page range if specified
        if page_range:
            try:
                page_numbers = PageRangeParser.parse(page_range, total_page_count)
                logger.info(
                    f"Processing {PageRangeParser.format_page_summary(page_numbers)} of {pdf_path.name}"
                )
            except Exception as e:
                logger.error(f"Invalid page range '{page_range}': {e}")
                raise ValueError(f"Invalid page range: {e}")
        else:
            page_numbers = list(range(1, total_page_count + 1))
            logger.info(f"Processing all {total_page_count} pages of {pdf_path.name}")

        # Initialize progress monitor
        progress_monitor = PageProgressMonitor(total_page_count, page_numbers)

        # Convert only specified pages to PIL Images
        if len(page_numbers) == total_page_count and page_numbers == list(
            range(1, total_page_count + 1)
        ):
            # Converting all pages - use simple approach
            progress_monitor.start_processing(1)
            images = convert_from_path(str(pdf_path), dpi=dpi, fmt="RGB", poppler_path=poppler_path)
            progress_monitor.finish_processing(total_page_count)
        else:
            # Converting specific pages - use page-by-page approach for progress monitoring
            images = []
            for page_num in page_numbers:
                progress_monitor.start_processing(page_num)
                try:
                    # Convert single page
                    page_images = convert_from_path(
                        str(pdf_path),
                        dpi=dpi,
                        fmt="RGB",
                        poppler_path=poppler_path,
                        first_page=page_num,
                        last_page=page_num,
                    )
                    if page_images:
                        images.extend(page_images)
                        progress_monitor.finish_processing(page_num, success=True)
                    else:
                        logger.warning(f"No image returned for page {page_num}")
                        progress_monitor.finish_processing(page_num, success=False)
                except Exception as e:
                    logger.error(f"Failed to convert page {page_num}: {e}")
                    progress_monitor.finish_processing(page_num, success=False)
                    raise

        # Convert images to data URIs
        data_uris = []
        processed_pages = len(images)

        for i, image in enumerate(images):
            # Convert PIL Image to base64 data URI
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG", quality=85)
            img_bytes = buffer.getvalue()

            # Create data URI
            base64_string = base64.b64encode(img_bytes).decode("utf-8")
            data_uri = f"data:image/jpeg;base64,{base64_string}"
            data_uris.append(data_uri)

            logger.debug(f"Converted page image {i + 1}/{processed_pages} to data URI")

        logger.info(
            f"Converted {processed_pages} pages from {pdf_path.name} (total pages: {total_page_count})"
        )
        return data_uris, total_page_count

    except Exception as e:
        logger.error(f"Failed to convert PDF {pdf_path} to data URIs: {e}")
        raise ValueError(f"PDF conversion failed: {e}")


# Load model from config when needed


class DocumentType(Enum):
    MARKDOWN = "markdown"
    DATASHEET_PDF = "datasheet_pdf"
    GENERIC_PDF = "generic_pdf"
    WORD_DOCUMENT = "word_document"
    POWERPOINT_PRESENTATION = "powerpoint_presentation"


class DocumentClassifier:
    """Classify documents to determine parsing strategy."""

    @staticmethod
    def classify(source: str | Path, is_datasheet_mode: bool = True) -> DocumentType:
        """Classify document type based on file extension and heuristics."""
        path = Path(source) if isinstance(source, Path) else Path(str(source))

        # Markdown and text files - no model call needed
        if path.suffix.lower() in {".md", ".markdown", ".txt"}:
            return DocumentType.MARKDOWN

        # PDF files - check if datasheet mode with additional heuristics
        if path.suffix.lower() == ".pdf":
            return DocumentClassifier._classify_pdf(path, is_datasheet_mode)

        # Office documents - Word and PowerPoint
        if path.suffix.lower() in {".docx", ".doc"}:
            return DocumentType.WORD_DOCUMENT

        if path.suffix.lower() in {".pptx", ".ppt"}:
            return DocumentType.POWERPOINT_PRESENTATION

        raise ValueError(f"Unsupported file type: {path.suffix}")

    @staticmethod
    def _classify_pdf(path: Path, is_datasheet_mode: bool) -> DocumentType:
        """Classify PDF based on filename patterns and mode."""
        filename_lower = path.name.lower()

        # Strong indicators for datasheets
        datasheet_indicators = [
            "datasheet",
            "ds.pdf",
            "spec",
            "specification",
            "product_brief",
            "technical_data",
            "sensor",
            "laser",
            "manual",
            "model",
            "part_number",
        ]

        # Strong indicators for generic documents
        generic_indicators = [
            "report",
            "paper",
            "article",
            "research",
            "white_paper",
            "guide",
            "tutorial",
            "documentation",
            "readme",
        ]

        # Check filename patterns
        has_datasheet_pattern = any(
            indicator in filename_lower for indicator in datasheet_indicators
        )
        has_generic_pattern = any(indicator in filename_lower for indicator in generic_indicators)

        # Decision logic
        if has_datasheet_pattern and not has_generic_pattern:
            logger.info(f"Detected datasheet pattern in filename: {path.name}")
            return DocumentType.DATASHEET_PDF
        if has_generic_pattern and not has_datasheet_pattern:
            logger.info(f"Detected generic document pattern in filename: {path.name}")
            return DocumentType.GENERIC_PDF
        if is_datasheet_mode:
            logger.info(f"Using datasheet mode for PDF: {path.name}")
            return DocumentType.DATASHEET_PDF
        logger.info(f"Using generic mode for PDF: {path.name}")
        return DocumentType.GENERIC_PDF

    @staticmethod
    def get_confidence(source: str | Path, doc_type: DocumentType) -> float:
        """Get confidence score for classification."""
        path = Path(source) if isinstance(source, Path) else Path(str(source))

        if doc_type == DocumentType.MARKDOWN:
            return 1.0  # Always confident about markdown files

        if doc_type == DocumentType.WORD_DOCUMENT:
            return 1.0  # Always confident about Word documents based on extension

        if doc_type == DocumentType.POWERPOINT_PRESENTATION:
            return 1.0  # Always confident about PowerPoint documents based on extension

        filename_lower = path.name.lower()

        if doc_type == DocumentType.DATASHEET_PDF:
            datasheet_indicators = [
                "datasheet",
                "ds.pdf",
                "spec",
                "specification",
                "product_brief",
                "technical_data",
                "sensor",
                "laser",
            ]
            matches = sum(1 for indicator in datasheet_indicators if indicator in filename_lower)
            return min(0.9, 0.5 + (matches * 0.2))  # 0.5-0.9 based on matches

        if doc_type == DocumentType.GENERIC_PDF:
            generic_indicators = [
                "report",
                "paper",
                "article",
                "research",
                "white_paper",
            ]
            matches = sum(1 for indicator in generic_indicators if indicator in filename_lower)
            return min(0.9, 0.5 + (matches * 0.2))  # 0.5-0.9 based on matches

        return 0.5  # Default medium confidence


async def parse_document(
    pdf_path: Path,
    doc_type: DocumentType,
    prompt_text: str,
    cache: CacheManager | None = None,
    config: PipelineConfig | None = None,
    page_range: str | None = None,
) -> tuple[str, list[tuple[str, str]], dict[str, Any]]:
    """Parse document based on type."""

    # Check cache first
    if cache:
        # Generate robust content-based hash for cache key
        try:
            # Read file content for hashing
            file_content = pdf_path.read_bytes()
            content_hash = hashlib.sha256(file_content).hexdigest()
        except Exception:
            # Fallback to path-based hash if content read fails
            content_hash = hashlib.sha256(str(pdf_path).encode()).hexdigest()

        # Create cache key that includes content + prompt + doc_type
        prompt_hash = hashlib.sha256(prompt_text.encode()).hexdigest()[:12]
        cache_key = f"{doc_type.value}_{content_hash[:12]}_{prompt_hash}"

        cached = cache.get(content_hash, cache_key)
        if cached:
            return cached["markdown"], cached["pairs"], cached["metadata"]

    if doc_type == DocumentType.MARKDOWN:
        # Direct read - no API call
        markdown = pdf_path.read_text(encoding="utf-8", errors="ignore")
        pairs = []  # No model/part pairs in markdown
        metadata = {
            "source_type": "markdown",
            "file_name": pdf_path.name,
            "file_size": pdf_path.stat().st_size,
            "content_length": len(markdown),
            "parse_method": "direct_read",
        }

    elif doc_type == DocumentType.DATASHEET_PDF:
        # Use special datasheet prompt with pair extraction
        markdown, pairs, _ = await vision_parse_datasheet(pdf_path, prompt_text, config, page_range)
        metadata = {
            "source_type": "datasheet_pdf",
            "extracted_pairs": len(pairs),
            "file_name": pdf_path.name,
            "file_size": pdf_path.stat().st_size,
            "content_length": len(markdown),
            "parse_method": "openai_vision",
        }

    elif doc_type == DocumentType.GENERIC_PDF:
        # Use generic prompt without pair extraction
        markdown, _, _ = await vision_parse_generic(pdf_path, prompt_text, config, page_range)
        pairs = []
        metadata = {
            "source_type": "generic_pdf",
            "file_name": pdf_path.name,
            "file_size": pdf_path.stat().st_size,
            "content_length": len(markdown),
            "parse_method": "openai_vision",
        }

    elif doc_type == DocumentType.WORD_DOCUMENT:
        # Parse Word document with python-docx
        markdown, pairs, _ = await parse_word_document(pdf_path, config)
        metadata = {
            "source_type": "word_document",
            "file_name": pdf_path.name,
            "file_size": pdf_path.stat().st_size,
            "content_length": len(markdown),
            "parse_method": "python_docx",
            "extracted_pairs": len(pairs),
        }

    elif doc_type == DocumentType.POWERPOINT_PRESENTATION:
        # Parse PowerPoint with python-pptx
        markdown, pairs, slide_metadata = await parse_powerpoint_document(pdf_path, config)
        metadata = {
            "source_type": "powerpoint_presentation",
            "file_name": pdf_path.name,
            "file_size": pdf_path.stat().st_size,
            "content_length": len(markdown),
            "parse_method": "python_pptx",
            "slide_count": slide_metadata.get("slide_count", 0),
        }

    # Cache result
    if cache:
        # Use the same content hash and cache key generated earlier
        try:
            file_content = pdf_path.read_bytes()
            content_hash = hashlib.sha256(file_content).hexdigest()
        except Exception:
            content_hash = hashlib.sha256(str(pdf_path).encode()).hexdigest()

        prompt_hash = hashlib.sha256(prompt_text.encode()).hexdigest()[:12]
        cache_key = f"{doc_type.value}_{content_hash[:12]}_{prompt_hash}"

        cache.put(
            content_hash,
            cache_key,
            {"markdown": markdown, "pairs": pairs, "metadata": metadata},
        )

    return markdown, pairs, metadata


async def vision_parse_datasheet(
    pdf: Path,
    parsing_prompt: str,
    config: PipelineConfig | None = None,
    page_range: str | None = None,
) -> tuple[str, list[tuple[str, str]], dict[str, Any]]:
    """Parse datasheet PDF with model/part number extraction."""
    client = create_vision_client(config)

    # Get model from config or use default
    model = config.openai.vision_model if config else "gpt-4.1"
    max_retries = config.openai.max_retries if config else 3

    # Enhanced prompt structure from notebook
    parts = [
        {
            "type": "input_text",
            "text": (
                f"{parsing_prompt}\n\n"
                "## ADDITIONAL INSTRUCTIONS\n"
                "Return **one Markdown document** with two clearly-separated sections:\n"
                "1. `Metadata:` keep exactly the JSON structure shown below and fill the "
                "`pairs` list you extracted (no extra keys).\n"
                "2. The **entire datasheet** translated into GitHub-flavoured Markdown, "
                "preserving all tables, headings, lists, line-breaks, and footnotes.\n\n"
                "Example top of output (do not include the ``` fences):\n"
                "Metadata: {\n"
                "    'pairs': [\n"
                "        ('PM10K+ DB-25 + USB', '2293937'),\n"
                "        ('PM10K+ RS-232', '2293938')\n"
                "    ]\n"
                "}\n\n"
                "---  ← leave one blank line, then start the document body ---\n"
            ),
        }
    ]

    # Add PDF pages as images (Responses API format)
    dpi = config.pdf.dpi if config and hasattr(config, "pdf") else 150
    data_uris, page_count = _pdf_to_data_uris(pdf, dpi=dpi, page_range=page_range)
    parts += [{"type": "input_image", "image_url": uri} for uri in data_uris]

    # Calculate timeout based on page count
    timeout_per_page = config.openai.timeout_per_page if config else 30
    timeout_base = config.openai.timeout_base if config else 60
    api_timeout = timeout_base + (page_count * timeout_per_page)
    logger.info(f"Using timeout of {api_timeout}s for {page_count} pages")

    # Make API call with retry using Responses API
    @enhanced_retry_api_call(max_attempts=max_retries, timeout=api_timeout, retry_type="vision")
    async def call_api():
        return client.responses.create(
            model=model,
            input=[{"role": "user", "content": parts}],
            temperature=0.0,
        )

    response = await call_api()
    md = response.output[0].content[0].text

    # Extract pairs from metadata block (handles multi-line JSON)
    try:
        if md.startswith("Metadata:"):
            # Find the end of metadata block (marked by "---" separator)
            if "\n---\n" in md:
                metadata_section, markdown_content = md.split("\n---\n", 1)
                # Extract JSON from metadata section
                json_text = metadata_section.replace("Metadata:", "").strip()
                # Handle single quotes in the response by converting to double quotes
                json_text = json_text.replace("'", '"')
                meta = json.loads(json_text)
                pairs = [tuple(p) for p in meta.get("pairs", [])]
                # Use content after the separator as markdown
                md = markdown_content.strip()
                logger.info(f"Successfully extracted {len(pairs)} pairs from metadata")
            else:
                # Fallback: try to parse first line only (backward compatibility)
                first_line, *rest = md.split("\n", 1)
                if first_line.startswith("Metadata:"):
                    metadata_text = first_line.replace("Metadata:", "").strip()
                    metadata_text = metadata_text.replace("'", '"')
                    meta = json.loads(metadata_text)
                    pairs = [tuple(p) for p in meta.get("pairs", [])]
                    md = "\n".join(rest) if rest else md
                else:
                    pairs = []
        else:
            pairs = []
    except Exception as e:
        logger.warning(f"Failed to extract pairs from metadata: {e}")
        logger.debug(f"Metadata section being parsed: {md[:200]}...")
        pairs = []

    # Return 3 values to match expected interface
    # Third value is metadata dict (empty for PDFs as metadata is in pairs)
    return md, pairs, {}


async def vision_parse_generic(
    pdf: Path,
    parsing_prompt: str,
    config: PipelineConfig | None = None,
    page_range: str | None = None,
) -> tuple[str, list[tuple[str, str]], dict[str, Any]]:
    """Parse generic PDF without pair extraction."""
    client = create_vision_client(config)

    # Get model from config or use default
    model = config.openai.vision_model if config else "gpt-4.1"
    max_retries = config.openai.max_retries if config else 3

    # Enhanced prompt for generic PDFs
    parts = [
        {
            "type": "input_text",
            "text": (
                parsing_prompt
                or "Extract all text from this document as GitHub-flavoured Markdown.\n\n"
                "## INSTRUCTIONS\n"
                "- Preserve all tables, headings, lists, and formatting\n"
                "- Maintain document structure and hierarchy\n"
                "- Include any technical specifications or data\n"
                "- Return **only** the Markdown content\n"
            ),
        }
    ]

    # Add PDF pages as images with configurable DPI
    dpi = config.pdf.dpi if config and hasattr(config, "pdf") else 150
    data_uris, page_count = _pdf_to_data_uris(pdf, dpi=dpi, page_range=page_range)
    parts += [{"type": "input_image", "image_url": uri} for uri in data_uris]

    # Calculate timeout based on page count
    timeout_per_page = config.openai.timeout_per_page if config else 30
    timeout_base = config.openai.timeout_base if config else 60
    api_timeout = timeout_base + (page_count * timeout_per_page)
    logger.info(f"Using timeout of {api_timeout}s for {page_count} pages")

    @enhanced_retry_api_call(max_attempts=max_retries, timeout=api_timeout, retry_type="vision")
    async def call_api():
        return client.responses.create(
            model=model,
            input=[{"role": "user", "content": parts}],
            temperature=0.0,
        )

    response = await call_api()
    # Return 3 values to match expected interface
    # Generic PDFs don't extract pairs, so return empty list and dict
    return response.output[0].content[0].text, [], {}


async def parse_word_document(
    word_path: Path, config: PipelineConfig | None = None
) -> tuple[str, list[tuple[str, str]], dict[str, Any]]:
    """Parse Word document with python-docx and convert to markdown."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx library is required for Word document parsing. Install with: pip install python-docx"
        )

    try:
        doc = Document(word_path)

        # Extract document metadata
        props = doc.core_properties
        logger.info(f"Parsing Word document: {word_path.name}")
        if props.title:
            logger.info(f"Document title: {props.title}")
        if props.author:
            logger.info(f"Document author: {props.author}")

        # Parse document structure
        sections = []

        # Process paragraphs and headings
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Determine if this is a heading
                style_name = paragraph.style.name
                if style_name.startswith("Heading"):
                    level = _get_heading_level(style_name)
                    sections.append(
                        {
                            "type": "heading",
                            "level": level,
                            "text": paragraph.text.strip(),
                        }
                    )
                else:
                    sections.append(
                        {
                            "type": "paragraph",
                            "text": paragraph.text.strip(),
                            "style": style_name,
                        }
                    )

        # Process tables
        for table in doc.tables:
            table_data = _extract_word_table(table)
            if table_data:
                sections.append({"type": "table", "data": table_data})

        # Convert to markdown
        markdown = _convert_word_sections_to_markdown(sections)

        # Extract model/part pairs for technical documents
        pairs = _extract_word_pairs(sections)

        # Create metadata
        metadata = {
            "source_type": "word_document",
            "file_name": word_path.name,
            "file_size": word_path.stat().st_size,
            "content_length": len(markdown),
            "doc_type": DocumentType.WORD_DOCUMENT.value,
            "section_count": len(sections),
            "pair_count": len(pairs),
            "title": props.title or "",
            "author": props.author or "",
            "created": props.created.isoformat() if props.created else None,
            "modified": props.modified.isoformat() if props.modified else None,
        }

        logger.info(
            f"Successfully parsed Word document with {len(sections)} sections and {len(pairs)} pairs"
        )
        return markdown, pairs, metadata

    except Exception as e:
        logger.error(f"Failed to parse Word document {word_path}: {e}")
        raise ValueError(f"Word document parsing failed: {e}")


async def parse_powerpoint_document(
    ppt_path: Path, config: PipelineConfig | None = None
) -> tuple[str, list[tuple[str, str]], dict[str, Any]]:
    """Parse PowerPoint presentation with python-pptx and convert to markdown."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx library is required for PowerPoint parsing. Install with: pip install python-pptx"
        )

    try:
        prs = Presentation(ppt_path)

        logger.info(f"Parsing PowerPoint presentation: {ppt_path.name}")
        logger.info(f"Slide count: {len(prs.slides)}")

        # Extract presentation metadata
        slide_contents = []

        for i, slide in enumerate(prs.slides):
            slide_data = _parse_slide_content(slide, i + 1)
            slide_contents.append(slide_data)

        # Convert slides to markdown
        markdown = _convert_slides_to_markdown(slide_contents)

        # PowerPoint presentations rarely have model/part pairs
        pairs = []

        # Create metadata
        metadata = {
            "source_type": "powerpoint_presentation",
            "file_name": ppt_path.name,
            "file_size": ppt_path.stat().st_size,
            "content_length": len(markdown),
            "doc_type": DocumentType.POWERPOINT_PRESENTATION.value,
            "slide_count": len(prs.slides),
            "has_speaker_notes": any(slide.get("notes") for slide in slide_contents),
        }

        logger.info(f"Successfully parsed PowerPoint with {len(prs.slides)} slides")
        return markdown, pairs, metadata

    except Exception as e:
        logger.error(f"Failed to parse PowerPoint presentation {ppt_path}: {e}")
        raise ValueError(f"PowerPoint parsing failed: {e}")


def _get_heading_level(style_name: str) -> int:
    """Extract heading level from Word style name."""
    try:
        if "Heading" in style_name:
            # Extract number from "Heading 1", "Heading 2", etc.
            parts = style_name.split()
            if len(parts) > 1 and parts[1].isdigit():
                return int(parts[1])
        return 1  # Default to H1 if can't determine
    except Exception:
        return 1


def _extract_word_table(table) -> list[list[str]]:
    """Extract table data from Word table."""
    try:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            if any(row_data):  # Only add non-empty rows
                table_data.append(row_data)
        return table_data
    except Exception as e:
        logger.warning(f"Failed to extract table data: {e}")
        return []


def _convert_word_sections_to_markdown(sections: list[dict]) -> str:
    """Convert Word document sections to markdown format."""
    markdown_parts = []

    for section in sections:
        if section["type"] == "heading":
            level = section["level"]
            heading_prefix = "#" * min(level, 6)  # Limit to H6
            markdown_parts.append(f"{heading_prefix} {section['text']}\n")

        elif section["type"] == "paragraph":
            markdown_parts.append(f"{section['text']}\n")

        elif section["type"] == "table":
            table_md = _convert_table_to_markdown(section["data"])
            markdown_parts.append(table_md)

    return "\n".join(markdown_parts)


def _convert_table_to_markdown(table_data: list[list[str]]) -> str:
    """Convert table data to markdown table format."""
    if not table_data:
        return ""

    markdown_lines = []

    # Header row
    if table_data:
        header = " | ".join(table_data[0])
        markdown_lines.append(f"| {header} |")

        # Separator row
        separator = " | ".join(["---"] * len(table_data[0]))
        markdown_lines.append(f"| {separator} |")

        # Data rows
        for row in table_data[1:]:
            # Pad row to match header length
            padded_row = row + [""] * (len(table_data[0]) - len(row))
            row_text = " | ".join(padded_row)
            markdown_lines.append(f"| {row_text} |")

    return "\n".join(markdown_lines) + "\n"


def _extract_word_pairs(sections: list[dict]) -> list[tuple[str, str]]:
    """Extract model/part number pairs from Word document sections."""
    pairs = []

    # Look for technical patterns in text
    import re

    pair_patterns = [
        r"([A-Z][A-Z0-9\-\+]+)\s*[:\-]\s*([0-9]{6,})",  # Model: partnumber
        r"Model\s*([A-Z][A-Z0-9\-\+]+).*?Part\s*(?:Number|No\.?)\s*([0-9]{6,})",  # Model X Part Number Y
        r"([A-Z]{2,}[0-9]+[A-Z\-\+]*)\s*\(([0-9]{6,})\)",  # MODEL(partnumber)
    ]

    for section in sections:
        if section["type"] in ["paragraph", "heading"]:
            text = section["text"]
            for pattern in pair_patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                for match in matches:
                    model, part = match
                    pairs.append((model.strip(), part.strip()))

    # Remove duplicates while preserving order
    seen = set()
    unique_pairs = []
    for pair in pairs:
        if pair not in seen:
            seen.add(pair)
            unique_pairs.append(pair)

    return unique_pairs


def _parse_slide_content(slide, slide_number: int) -> dict[str, Any]:
    """Parse content from a single PowerPoint slide."""
    content = {
        "slide_number": slide_number,
        "title": "",
        "content": [],
        "tables": [],
        "notes": "",
    }

    # Extract text from shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Check if this is a title placeholder
            if (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format
                and shape.placeholder_format.type == 1
            ):  # Title placeholder
                content["title"] = shape.text.strip()
            else:
                content["content"].append(shape.text.strip())

        # Handle tables in slides
        if hasattr(shape, "table") and shape.has_table:
            table_data = _extract_ppt_table(shape.table)
            if table_data:
                content["tables"].append(table_data)

    # Extract speaker notes
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            content["notes"] = notes_text

    return content


def _extract_ppt_table(table) -> list[list[str]]:
    """Extract table data from PowerPoint table."""
    try:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            if any(row_data):  # Only add non-empty rows
                table_data.append(row_data)
        return table_data
    except Exception as e:
        logger.warning(f"Failed to extract PowerPoint table data: {e}")
        return []


def _convert_slides_to_markdown(slide_contents: list[dict]) -> str:
    """Convert PowerPoint slides to markdown format."""
    markdown_parts = []

    for slide in slide_contents:
        slide_num = slide["slide_number"]
        title = slide["title"]
        content = slide["content"]
        tables = slide["tables"]
        notes = slide["notes"]

        # Slide header
        if title:
            markdown_parts.append(f"# Slide {slide_num}: {title}\n")
        else:
            markdown_parts.append(f"# Slide {slide_num}\n")

        # Slide content
        for text in content:
            # Handle bullet points and lists
            if text.startswith(("•", "-")):
                markdown_parts.append(f"{text}\n")
            else:
                markdown_parts.append(f"{text}\n")

        # Tables
        for table_data in tables:
            table_md = _convert_table_to_markdown(table_data)
            markdown_parts.append(table_md)

        # Speaker notes
        if notes:
            markdown_parts.append(f"\n**Speaker Notes:** {notes}\n")

        markdown_parts.append("\n---\n")  # Slide separator

    return "\n".join(markdown_parts)
